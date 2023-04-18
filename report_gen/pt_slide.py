"""
The RVA Outbrief slide generator takes a presentation template and
information from the RVA pen testing django app database to generate
an RVA Outbrief slide deck in pptx format
"""
# Risk & Vulnerability Assessment Reporting Engine

# Copyright 2022 The Risk & Vulnerability Reporting Engine Contributors, All Rights Reserved.
# (see Contributors.txt for a full list of Contributors)

# SPDX-License-Identifier: BSD-3-Clause

# Please see additional acknowledgments (including references to third party source code, object code, documentation and other files) in the license.txt file or contact permission@sei.cmu.edu for full terms.

# Created, in part, with funding and support from the United States Government. (see Acknowledgments file).

# DM22-1011

from os import terminal_size
import sys
import os.path
import argparse
import datetime
from datetime import timezone
from dateutil.relativedelta import relativedelta
import report_gen.utilities.rt_parser as rt

try:
    import pptx
except ImportError:
    print("Must have python-pptx library installed")
    sys.exit(1)

if pptx.__version__ < "0.6.18":
    print("Must have python-pptx version 0.6.18 or greater")
    sys.exit(1)

import report_gen.utilities.assessment_facts as af
import report_gen.utilities.img_util as iu
from pptx.util import Inches, Pt

from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Constants
tstamp = str(datetime.datetime.now().strftime("%Y%m%d_%H.%M.%S"))
pr_tstamp = str(datetime.datetime.now().strftime("%B %d, %Y"))

red = RGBColor(255, 0, 0)
green = RGBColor(70, 200, 0)  # pure green is hard to read
blue = RGBColor(0, 82, 136)
gray = RGBColor(90, 91, 92)
white = RGBColor(255, 255, 255)
blue = RGBColor(3, 82, 136)
crit = RGBColor(255, 116, 113)
high = RGBColor(252, 191, 143)
med = RGBColor(255, 222, 89)
low = RGBColor(131, 224, 142)
info = RGBColor(79, 175, 227)
LightStyle2Accent4 = "{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}"
MediumStyle4Accent4 = "{C4B1156A-380E-4F78-BDF5-A606A8083BF9}"


def iter_cells(table):
    """Walks through the cells of a table.

    Args:
        table (pptx table object): The table that will be iterated through.

    Yields:
        [Generator]: Returns a generator object that contains each cell inside of the table.
    """
    for row in table.rows:
        for cell in row.cells:
            yield cell


def table_adjust_fontsize(table, sz):
    """Adjust the text size in all of the cells of a table.

    Args:
        table ([type]): [description]
        sz ([type]): [description]
    """
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(sz)


def cell_text(t, r, c, s, alignment="n", color=None, size=None):
    """Helper function to assist in setting various options for cells in a table.

    Args:
        t (pptx table): The table that contains the cell being changed.
        r (int): The row of the cell.
        c (int): The column of the cell.
        s (str): The text inside of the cell.
        alignment (str, optional): Sets the alignment of the text inside of the cell. Defaults to 'n'.
        color (RGBColor, optional): The color of the text to be set. Defaults to None.
        size (int, optional): The size to set the text in the cell to. Defaults to None.
    """
    t.cell(r, c).text = s
    if color is not None:
        t.cell(r, c).text_frame.paragraphs[0].font.color.rgb = color

    if size is not None:
        t.cell(r, c).text_frame.paragraphs[0].font.size = size

    if alignment == "n":
        # no alignment specified
        return
    elif alignment == "c":
        alignment = PP_ALIGN.CENTER
    elif alignment == "l":
        alignment = PP_ALIGN.LEFT
    elif alignment == "j":
        alignment = PP_ALIGN.JUSTIFY
    else:
        alignment = PP_ALIGN.CENTER

    t.cell(r, c).text_frame.paragraphs[0].alignment = alignment


def blocked(s):
    """Returns proper strings and color of phishing table entries.

    Args:
        s (str): The short string of the phishing cell.

    Returns:
        (str, RGBColor): Returns a tuble containing the proper string and color for the text.
    """
    if s == "B":
        return ("Blocked", green)
    elif s == "N":
        return ("Not Blocked", red)
    else:
        return ("Not Set", None)


def add_hyperlink(para, link_text, link_address):
    """Adds a hyperlinked text inside of the slide.

    Args:
        para (pptx paragraph): The paragarph that the hyperlink will be inserted at.
        link_text (str): The text that holds the link.
        link_address (str): Where the hyperlink redirects to.
    """
    if link_text is None:
        link_text = link_address

    run = para.add_run()
    run.text = link_text
    hlink = run.hyperlink
    hlink.address = link_address


def add_tag(shapes, tag_type, label):

    if tag_type == "severity":
        tag = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.64), Inches(6.69), Inches(1.83), Inches(0.42))
        tag.fill.solid()

        if label == "Critical":
            tag.fill.fore_color.rgb = crit
            tag.line.color.rgb = crit
        elif label == "High":
            tag.fill.fore_color.rgb = high
            tag.line.color.rgb = high
        elif label == "Medium":
            tag.fill.fore_color.rgb = med
            tag.line.color.rgb = med
        elif label == "Low":
            tag.fill.fore_color.rgb = low
            tag.line.color.rgb = low
        else:
            tag.fill.fore_color.rgb = info
            tag.line.color.rgb = info

    elif tag_type == "mitigation":
        tag = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.09), Inches(6.69), Inches(1.83), Inches(0.42))
        tag.fill.solid()

        if label == "Mitigated":
            tag.fill.fore_color.rgb = low
            tag.line.color.rgb = low
        else:
            tag.fill.fore_color.rgb = crit
            tag.line.color.rgb = crit

    else:
        tag = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.53), Inches(6.69), Inches(1.83), Inches(0.42))
        tag.fill.solid()
        tag.fill.fore_color.rgb = crit
        tag.line.color.rgb = crit

    shadow = tag.shadow
    shadow.inherit = False
    tag_text = tag.text_frame.paragraphs[0]
    tag_text.text = label
    tag_text.font.size = Pt(16)
    tag_text.font.color.rgb = gray
    tag_text.runs[0].font.bold = True


def add_section_title(prs, title):
    """Adds a section slide into the powerpoint.

    Args:
        prs (pptx presentation): The powerpoint presentation.
        title (str): Name of the section that will be inserted in the slide.
    """
    section_slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(section_slide_layout)
    slide.shapes.title.text = title


def auto_size(shape):
    """Attempts to autosize the text inside of a textbox.

    Args:
        shape (powerpoint shape): The shape in which the text is being resized.
    """
    total_runs = 0
    font_size = 26

    for paragraph in shape.paragraphs:
        total_runs += len(paragraph.runs)

    if total_runs > 15:
        # don't font to go below 10 points.
        # user can always switch to 2 column mode manually.
        font_size = max(10, font_size - (total_runs - 5))

    for paragraph in shape.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def insert_title_slide(prs, report_type, rva_info, draft):
    """Inserts the title slide into the presentation.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
        rva_info (List[Dict]): The Json data from the engagement.
        draft (bool): Boolean to mark if the presentation is a draft or not.
    """
    open_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(open_slide_layout)

    subtitle = slide.placeholders[10]
    date = slide.placeholders[11]

    customer = af.get_db_info(
        rva_info, "engagementmeta.fields.customer_long_name", "Stakeholder Long Name"
    )
    customer_initials = af.get_db_info(
        rva_info, "engagementmeta.fields.customer_initials", "Stakeholder Initials"
    )

    subtitle.text = customer + " (" + customer_initials + ")"
    date.text = pr_tstamp


def insert_notice_slide(prs, report_type):
    """Insert the notice slide into the powerpoint.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
    """
    notice_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(notice_slide_layout)


def insert_agenda_slide(prs, report_type):
    """Generates the agenda slide. Puts in specific agenda depending on the report.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
    """
    agenda_slide_layout = prs.slide_layouts[10]
    slide = prs.slides.add_slide(agenda_slide_layout)


def insert_logistics(prs, report_type, rva_info):
    """Inserts the slide with details of the timeframe from the engagement.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
        rva_info (List[Dict]): The Json data from the engagement.
    """
    title_only = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only)
    shapes = slide.shapes
    shapes.title.text = "LOGISTICS"

    if report_type != "RPT":
        rows = 3
    else:
        rows = 2

    # date table
    cols = 2
    left = Inches(0.65)
    top = Inches(1.58)
    width = Inches(8.71)
    height = Inches(1.2)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    tbl = table._graphic_frame._element.graphic.graphicData.tbl
    tbl[0][-1].text = LightStyle2Accent4

    # set column widths
    table.columns[0].width = Inches(4.06)
    table.columns[1].width = Inches(4.65)

    # write body cells
    ext_start_date = af.get_db_info(
        rva_info, "engagementmeta.fields.ext_start_date", "External Start Date"
    )
    ext_start_formatted = datetime.datetime.strptime(ext_start_date, '%Y-%m-%d').strftime('%b %-d, %Y')

    ext_end_date = af.get_db_info(
        rva_info, "engagementmeta.fields.ext_end_date", "External End Date"
    )
    ext_end_formatted = datetime.datetime.strptime(ext_end_date, '%Y-%m-%d').strftime('%b %-d, %Y')

    cell_text(table, 0, 0, "Activity", "c", white)
    cell_text(table, 0, 1, "Dates", "c", white)

    cell_text(table, 1, 0, "External Testing", color=gray)
    cell_text(table, 1, 1, str(ext_start_formatted) + " to " + str(ext_end_formatted), color=gray)

    if report_type != "RPT":
        int_start_date = af.get_db_info(
            rva_info, "engagementmeta.fields.int_start_date", "Internal Start Date"
        )
        int_start_formatted = datetime.datetime.strptime(int_start_date, '%Y-%m-%d').strftime('%b %-d, %Y')
        int_end_date = af.get_db_info(
            rva_info, "engagementmeta.fields.int_end_date", "Internal End Date"
        )
        int_end_formatted = datetime.datetime.strptime(int_end_date, '%Y-%m-%d').strftime('%b %-d, %Y')

        cell_text(table, 2, 0, "Internal Testing", color=gray)
        cell_text(table, 2, 1, str(int_start_formatted) + " to " + str(int_end_formatted), color=gray)
    
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = blue
    table.cell(0, 1).fill.solid()
    table.cell(0, 1).fill.fore_color.rgb = blue
    table.cell(1, 0).fill.solid()
    table.cell(1, 0).fill.fore_color.rgb = white
    table.cell(1, 1).fill.solid()
    table.cell(1, 1).fill.fore_color.rgb = white
    table.cell(2, 0).fill.solid()
    table.cell(2, 0).fill.fore_color.rgb = white
    table.cell(2, 1).fill.solid()
    table.cell(2, 1).fill.fore_color.rgb = white

    table.cell(1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(2, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(2, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # poc table
    rows = 2
    cols = 2
    left = Inches(0.65)
    top = Inches(3.12)
    width = Inches(8.71)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    tbl = table._graphic_frame._element.graphic.graphicData.tbl
    tbl[0][-1].text = LightStyle2Accent4

    # set column widths
    table.columns[0].width = Inches(4.06)
    table.columns[1].width = Inches(4.65)

    # write body cells
    poc_name = af.get_db_info(
        rva_info, "engagementmeta.fields.customer_POC_name", "Customer Name"
    )
    poc_email = af.get_db_info(
        rva_info, "engagementmeta.fields.customer_POC_email", "Customer Email"
    )

    cell_text(table, 0, 0, "Point of Contact", "c", color=white)
    table.cell(0, 0).merge(table.cell(0, 1))

    cell_text(table, 1, 0, poc_name, color=gray)
    cell_text(table, 1, 1, poc_email, color=gray)

    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = blue
    table.cell(1, 0).fill.solid()
    table.cell(1, 0).fill.fore_color.rgb = white
    table.cell(1, 1).fill.solid()
    table.cell(1, 1).fill.fore_color.rgb = white

    table.cell(1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # rva team table
    rows = 5
    cols = 2
    left = Inches(0.65)
    top = Inches(4.25)
    width = Inches(8.71)
    height = Inches(2)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    tbl = table._graphic_frame._element.graphic.graphicData.tbl
    tbl[0][-1].text = LightStyle2Accent4

    # set column widths
    table.columns[0].width = Inches(4.06)
    table.columns[1].width = Inches(4.65)

    # write body cells
    fed_name = af.get_db_info(
        rva_info, "engagementmeta.fields.team_lead_name", "Fed Lead Name"
    )

    cell_text(table, 0, 0, report_type + " Team", "c", color=white)
    table.cell(0, 0).merge(table.cell(0, 1))

    cell_text(table, 1, 0, "Federal Lead", color=gray)
    cell_text(table, 1, 1, fed_name, color=gray)
    cell_text(table, 2, 0, "Technical Lead", color=gray)
    cell_text(table, 2, 1, "<Tech Lead Name>", color=gray)
    cell_text(table, 3, 0, "Operator", color=gray)
    cell_text(table, 3, 1, "<Operator Name>", color=gray)
    cell_text(table, 4, 0, "Operator", color=gray)
    cell_text(table, 4, 1, "<Operator Name>", color=gray)

    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = blue

    for i in range(1, 5):
        table.cell(i, 0).fill.solid()
        table.cell(i, 0).fill.fore_color.rgb = white
        table.cell(i, 1).fill.solid()
        table.cell(i, 1).fill.fore_color.rgb = white

        table.cell(i, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        table.cell(i, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def insert_scope_slide(prs, report_type, rva_info, ip_ext_scan, ip_ext_disc, ip_int_scan, ip_int_disc):
    """Generates the slide containing the scope of the engagement.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
        rva_info (List[Dict]): The Json data from the engagement.
        ip_ext (int): The number of external IP addresses scanned.
        ip_int (int): The number of internal IP addresses scanned.
    """
    # ---- add scope and limitations slide
    scope_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(scope_slide_layout)
    scope_placeholder = slide.placeholders[1].text_frame

    title = slide.shapes.title
    title.text = "SCOPE"

    # external bullet
    paragraph = scope_placeholder.paragraphs[0]
    paragraph.level = 0
    run = paragraph.add_run()
    run.text = "External Scope"
    paragraph.font.color.rgb = gray

    # external information
    paragraph = scope_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = str(ip_ext_disc) + " active hosts out of " + str(ip_ext_scan) + " scanned hosts"
    paragraph.font.color.rgb = gray

    # internal bullet
    if report_type != "RPT":
        paragraph = scope_placeholder.add_paragraph()
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = "Internal Scope"
        paragraph.font.color.rgb = gray

        # internal information
        paragraph = scope_placeholder.add_paragraph()
        paragraph.level = 1
        run = paragraph.add_run()
        run.text = str(ip_int_disc) + " active hosts out of " + str(ip_int_scan) + " scanned hosts"
        paragraph.font.color.rgb = gray

    # testing limitations bullet
    paragraph = scope_placeholder.add_paragraph()
    paragraph.level = 0
    run = paragraph.add_run()
    run.text = "Testing Limitations"
    paragraph.font.color.rgb = gray

    # limitations
    paragraph = scope_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Short timeframe - overcome by working with {} staff".format(
        af.get_db_info(rva_info, "engagementmeta.fields.customer_initials", "CI")
    )
    paragraph.font.color.rgb = gray

    paragraph = scope_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Testing assumes in-scope systems are a fair representation of all production systems"
    paragraph.font.color.rgb = gray


def insert_goals_slide(prs, report_type):
    """Generates the goals of the assessment slide.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
    """

    goals_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(goals_slide_layout)


#def generate_narrative_data(db):
    """
    Generates and returns all of the data needed to populate the narrative section of the report.

    Args:
        db (List of dict(s)): The JSON output of the assessment facts.

    Returns:
        List: List with the needed narrative information.
    """
    """
    narrative = af.model_gen(db, "ptportal.narrative")
    # Narrative is sorted by PK to get all of the sections sorted based on their main narrative.
    sorted_narrative = sorted(narrative, key=lambda i: i["pk"])
    screenshots = af.model_gen(db, "ptportal.toolscreenshot")

    for s in screenshots:
        n = sorted_narrative[s.get("fields").get("narrative") - 1]
        # Screenshots are seperated from the main narrative so have to match up narratives with their respective screenshots.
        if "Screenshots" not in n:
            n["Screenshots"] = []
        n["Screenshots"].append(s)
    return sorted_narrative
    """

#def insert_narrative_slide(prs, report_type, rva_info, media_path):
    """Generates the narrative slides and inserts each step of the narrative.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
        rva_info (List[Dict]): The Json data from the engagement.
        media_path (str): The media file path.
    """
    """
    add_section_title(prs, "Narrative")

    narrative_data = generate_narrative_data(rva_info)
    order = ["Phishing", "External", "Internal"]

    for o in order:
        add_section_title(prs, o)

        narrative = [x for x in narrative_data if x.get("fields").get("type") == o]
        for n in narrative:
            if "Screenshots" not in n and n["fields"]["tool_output"] == "":
                continue

            if "Screenshots" in n:
                for s in n["Screenshots"]:
                    screenshot = prs.slides.add_slide(prs.slide_layouts[9])

                    screenshot.shapes.title.text = n["fields"]["name"]

                    for shape in screenshot.shapes:
                        if hasattr(shape, "text"):
                            if shape.text == "":
                                shape.text = s["fields"]["caption"]

                    sfile = media_path + s["fields"]["file"]
                    (x, y, w, h) = iu.get_screenshot_dimensions(sfile)
                    screenshot.shapes.add_picture(
                        sfile, Inches(x), Inches(y), width=Inches(w), height=Inches(h)
                    )

                    notes_slide = screenshot.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = (
                        "Image Caption\n"
                        + s["fields"]["caption"]
                        + "\n\nScreenshot Description\n"
                        + n["fields"]["tool_output_description"]
                    )

            if n["fields"]["tool_output"] != "":
                code_slide = prs.slides.add_slide(prs.slide_layouts[14])
                code_slide.shapes.title.text = n["fields"]["name"]

                code_paragraph = code_slide.placeholders[1].text_frame
                parser = rt.RichTextParser(code_paragraph)
                parser.feed(n["fields"]["tool_output"])
                parser.emit_pptx()

                for shape in code_slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text == "":
                            shape.text = n["fields"]["tool_output_description"]
    """
    """
    # ---- HVA attack overview results slides
    if report_type == "HVA":
        # ---- External Assessment slide
        hva_attack_overview(prs, rva_info,
                            "External Assessment",
                            "report.fields.HVA_external_scenario",
                            "External Assessment Scenario")

        # ---- Phishing Assessment slide
        hva_attack_overview(prs, rva_info,
                            "Phishing Assessment",
                            "report.fields.HVA_phishing_scenario",
                            "Phishing Assessment Scenario")

        # ---- Web Application Assessment slide
        hva_attack_overview(prs, rva_info,
                            "Web Application Assessment",
                            "report.fields.HVA_web_application_scenario",
                            "Web Application Assessment Scenario")

        # ---- Internal Assessment slide
        hva_attack_overview(prs, rva_info,
                            "Internal Assessment",
                            "report.fields.HVA_internal_scenario",
                            "Internal Assessment Scenario")

        # ---- Internal Threat Emulation slide
        hva_attack_overview(prs, rva_info,
                            "Internal Threat Emulation",
                            "report.fields.HVA_ITE_scenario",
                            "Internal Threat Emulation Scenario")

        # ---- Data Exfiltration slide
        hva_attack_overview(prs, rva_info,
                            "Data Exfiltration",
                            "report.fields.HVA_data_exfiltration_scenario",
                            "Data Exfiltration Scenario")
    """


def insert_OSINT_slide(prs, report_type, rva_info):
    """Generate the OSINT slides if the report is a RPT.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
        rva_info (List[Dict]): The Json data from the engagement.
    """
    osint_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(osint_slide_layout)

    rpt_info = af.get_db_info(rva_info, "report.fields", "keyNA")

    email_id = str(rpt_info["emails_identified"])
    email_brch = str(rpt_info["emails_breached"])
    creds_identified = str(rpt_info["credentials_identified"])
    creds_validated = str(rpt_info["credentials_validated"])

    if email_id == "0":
        email_id = "<NOT SET>"
    if email_brch == "0":
        email_brch = "<NOT SET>"
    if creds_identified == "0":
        creds_identified = "<NOT SET>"
    if creds_validated == "0":
        creds_validated = "<NOT SET>"

    title = slide.shapes.title
    agenda_placeholder = slide.placeholders[1].text_frame
    title.text = "Open Source Information Gathering"
    paragraph = agenda_placeholder.paragraphs[0]
    # Title is two lines at font size 40, add_paragraph() adds space from title to stop font overlap
    paragraph = agenda_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = email_id + " emails were scraped from various Internet sources."

    paragraph = agenda_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = (
        email_brch
        + " scraped emails were identified as existing in previous data breaches (according to HaveIBeenPwned database)"
    )

    paragraph = agenda_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = (
        creds_identified
        + " sets of credentials (emails and passwords) identified in the wild."
    )

    paragraph = agenda_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = creds_validated + " sets of credentials were validated."


def insert_findings_slides(prs, report_type, rva_info, ss_info, media_path):
    """Generates the slides for each finding in the assessment.

    Args:
        prs (pptx presentation): The Powerpoint presentation.
        report_type (str): The type of report being generated. I.e. RVA, RPT, or HVA.
        rva_info (List[Dict]): The Json data from the engagement.
        ss_info (list): The informatoin about all of the findings screenshots.
        media_path (str): The media file location.
    """

    # ---- add Findings Section Slide
    title_only = prs.slide_layouts[11]
    slide = prs.slides.add_slide(title_only)
    slide.placeholders[13].text = "FINDINGS"

    # ---- add Findings Severity Classification
    severity_rating_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(severity_rating_layout)

    # add Risk Score slide
    risk_score_layout = prs.slide_layouts[13]
    slide = prs.slides.add_slide(risk_score_layout)

    total_risk_score = 0
    mitigated_risk_score = 0

    findings = []

    for cnt, finding in enumerate(af.model_gen(rva_info, 'ptportal.uploadedfinding')):
        ele = finding['fields']
        if ele['KEV']:
            f_data = {"pk": finding['pk'], "name": ele['uploaded_finding_name'], "severity": ele['severity'], "mitigation": ele['mitigation'], "risk_score": ele['risk_score'], "kev": True}
        else:
            f_data = {"pk": finding['pk'], "name": ele['uploaded_finding_name'], "severity": ele['severity'], "mitigation": ele['mitigation'], "risk_score": ele['risk_score'], "kev": False}

        findings.append(f_data)

    for f in findings:
        total_risk_score += f['risk_score']
        if not f['mitigation']:
            mitigated_risk_score += f['risk_score']

    slide.placeholders[11].text = str(total_risk_score)
    slide.placeholders[12].text = str(mitigated_risk_score)

    risk_chart = media_path + 'charts/riskchart.png'

    slide.placeholders[13].insert_picture(risk_chart)

    order = {"Critical": 0, "High": 1, "Medium": 2, "Low": 3, "Informational": 4}
    findings_list = sorted(findings, key=lambda s: order[s['severity']])
    chunks = [findings_list[x:x+10] for x in range(0, len(findings_list), 10)]

    for cnt, chunk in enumerate(chunks):
        findings_summary_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(findings_summary_layout)
        shapes = slide.shapes
        shapes.title.text = "FINDINGS SUMMARY"

        # findings summary table
        rows = len(chunk) + 1
        cols = 3
        left = Inches(0.65)
        top = Inches(1.5)
        width = Inches(8.71)
        height = Inches(rows * 0.4)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        tbl = table._graphic_frame._element.graphic.graphicData.tbl
        tbl[0][-1].text = LightStyle2Accent4

        # set column widths
        table.columns[0].width = Inches(0.6)
        table.columns[1].width = Inches(6.25)
        table.columns[2].width = Inches(1.85)

        cell_text(table, 0, 0, "", "c", white)
        cell_text(table, 0, 1, "Finding Name", "c", white)
        cell_text(table, 0, 2, "Severity", "c", white)

        for i in range(0, 3):
            table.cell(0, i).fill.solid()
            table.cell(0, i).fill.fore_color.rgb = blue
            table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(0, i).vertical_anchor = MSO_ANCHOR.MIDDLE

        for count, finding in enumerate(chunk):
            f_id = (cnt * 10) + (count + 1)
            cell_text(table, count + 1, 0, str(f_id), color=gray)
            cell_text(table, count + 1, 1, finding['name'], color=gray)
            sev_cell = table.cell(count + 1, 2).text_frame.paragraphs[0]
            sev_color = sev_cell.add_run()
            sev_color.text =  ("• ")

            if finding['severity'] == "Critical":
                sev_color.font.color.rgb = crit
            elif finding['severity'] == "High":
                sev_color.font.color.rgb = high
            elif finding['severity'] == "Medium":
                sev_color.font.color.rgb = med
            elif finding['severity'] == "Low":
                sev_color.font.color.rgb = low
            else:
                sev_color.font.color.rgb = info

            severity = sev_cell.add_run()
            severity.text = finding['severity']
            severity.font.color.rgb = gray
            table.cell(count + 1, 0).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(count + 1, 2).text_frame.paragraphs[0].runs[0].font.bold = True
            table.cell(count + 1, 2).text_frame.paragraphs[0].runs[1].font.bold = True
            table.cell(count + 1, 2).text_frame.paragraphs[0].runs[1].font.size = Pt(11)

            for i in range(0, 3):
                table.cell(count + 1, i).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                table.cell(count + 1, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                table.cell(count + 1, i).vertical_anchor = MSO_ANCHOR.MIDDLE



    for cnt, finding in enumerate(findings_list):

        fpk = finding['pk']
        name = finding['name']
        severity = finding['severity']
        mitigation = finding['mitigation']
        kev = finding['kev']

        screenshots = af.find_screenshots(ss_info, finding['pk'])

        if len(screenshots) == 0:
            findings_layout = prs.slide_layouts[14]
            slide = prs.slides.add_slide(findings_layout)
            shapes = slide.shapes
            shapes.title.text = finding['name']

            add_tag(shapes, "severity", severity)

            if mitigation:
                add_tag(shapes, "mitigation", "Mitigated")
            else:
                add_tag(shapes, "mitigation", "Not Mitigated")

            if kev:
                add_tag(shapes, "kev", "KEV")

        else:
            for screenshot in screenshots:
                ssf = screenshot['fields']

                findings_layout = prs.slide_layouts[14]
                slide = prs.slides.add_slide(findings_layout)
                shapes = slide.shapes
                shapes.title.text = finding['name']

                sfile = media_path + ssf['file']
                (x, y, w, h) = iu.get_screenshot_dimensions(sfile, "finding")

                shapes.add_picture(sfile, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
                slide.placeholders[12].text = ssf['caption']

                add_tag(shapes, "severity", severity)

                if mitigation:
                    add_tag(shapes, "mitigation", "Mitigated")
                else:
                    add_tag(shapes, "mitigation", "Not Mitigated")

                if kev:
                    add_tag(shapes, "kev", "KEV")

    kevs = []

    for cnt, kev in enumerate(af.model_gen(rva_info, 'ptportal.kev')):
        ele = kev['fields']

        if ele['found']:
            found_kev = {"name": ele['vulnerability_name'], "cve": ele['cve_id']}
            kevs.append(found_kev)

    kevs_list = sorted(kevs, key=lambda k: k['cve'])
    kev_chunks = [kevs_list[x:x+10] for x in range(0, len(kevs_list), 10)]

    for cnt, chunk in enumerate(kev_chunks):
        kev_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(kev_layout)
        shapes = slide.shapes
        shapes.title.text = "KEV SUMMARY"

        # kev summary table
        rows = len(chunk) + 1
        cols = 2
        left = Inches(0.65)
        top = Inches(1.5)
        width = Inches(8.71)
        height = Inches(rows * 0.4)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        tbl = table._graphic_frame._element.graphic.graphicData.tbl
        tbl[0][-1].text = LightStyle2Accent4

        # set column widths
        table.columns[0].width = Inches(0.6)
        table.columns[1].width = Inches(8.1)

        cell_text(table, 0, 0, "", "c", white)
        cell_text(table, 0, 1, "Known Exploited Vulnerabilities (" + str(len(kevs_list)) + ")", "c", white)

        for i in range(0, 2):
            table.cell(0, i).fill.solid()
            table.cell(0, i).fill.fore_color.rgb = blue
            table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(0, i).vertical_anchor = MSO_ANCHOR.MIDDLE


        for count, kev in enumerate(chunk):
            k_id = (cnt * 10) + (count + 1)
            cell_text(table, count + 1, 0, str(k_id), color=gray)
            cell_text(table, count + 1, 1, kev['cve'] + ": " + kev['name'], color=gray)

            table.cell(count + 1, 0).text_frame.paragraphs[0].runs[0].font.bold = True

            for i in range(0, 2):
                table.cell(count + 1, i).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                table.cell(count + 1, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                table.cell(count + 1, i).vertical_anchor = MSO_ANCHOR.MIDDLE


def insert_services_slides(prs, rva_info):

    # insert ransomware slide
    if af.get_db_info(rva_info, 'ransomwarescenarios', 'NA') is not None or af.get_db_info(rva_info, 'ransomware', 'NA') is not None:
        ransomware_slide_layout = prs.slide_layouts[4]
        slide = prs.slides.add_slide(ransomware_slide_layout)
        ransomware_placeholder = slide.placeholders[1].text_frame

        title = slide.shapes.title
        title.text = "RANSOMWARE SUSCEPTIBILITY"

        vuln = 0
        total = 0
        if af.get_db_info(rva_info, 'ransomwarescenarios.fields.vuln', 'NA') != '<not set: NA>':
            vuln = int(af.get_db_info(rva_info, 'ransomwarescenarios.fields.vuln', 'NA'))
        if af.get_db_info(rva_info, 'ransomwarescenarios.fields.total', 'NA') != '<not set: NA>':
            total = int(af.get_db_info(rva_info, 'ransomwarescenarios.fields.total', 'NA'))

        paragraph = ransomware_placeholder.paragraphs[0]
        paragraph.level = 1
        run = paragraph.add_run()
        run.text = f"During ransomware simulation, the CISA team found that endpoints are vulnerable to {vuln} out of the {total} ransomware scenarios tested."
        paragraph.font.color.rgb = gray

        ransomware_results = []

        for item in af.model_gen(rva_info, 'ptportal.ransomware'):
            ele = item['fields']
            if not ele['disabled']:
                if "detected by security software" in ele['description']:
                    if ele['trigger'] == "Y":
                        if ele['time_start'] and ele['time_end']:
                            start = datetime.datetime.fromisoformat(ele['time_start'][:-1]).astimezone(timezone.utc)
                            end = datetime.datetime.fromisoformat(ele['time_end'][:-1]).astimezone(timezone.utc)
                            difference = relativedelta(end, start)

                            if difference.days == 1:
                                days = "1 day "
                            else:
                                days = str(difference.days) + " days "
                            if difference.hours == 1:
                                hours = "1 hour "
                            else:
                                hours = str(difference.hours) + " hours "
                            if difference.minutes == 1:
                                minutes = "and 1 minute."
                            else:
                                minutes = "and " + str(difference.minutes) + " minutes."
                        else:
                            days = "0 days "
                            hours = "0 hours "
                            minutes = "and 0 minutes."
                        ransomware_results.append("Ransomware activity was detected by security software within " + days + hours + minutes)
                    else:
                        ransomware_results.append("Ransomware activity was not detected by security software.")

                if "prevented by security software" in ele['description']:
                    if ele['trigger'] == "Y":
                        if ele['time_start'] and ele['time_end']:
                            start = datetime.datetime.fromisoformat(ele['time_start'][:-1]).astimezone(timezone.utc)
                            end = datetime.datetime.fromisoformat(ele['time_end'][:-1]).astimezone(timezone.utc)
                            difference = relativedelta(end, start)

                            if difference.days == 1:
                                days = "1 day "
                            else:
                                days = str(difference.days) + " days "
                            if difference.hours == 1:
                                hours = "1 hour "
                            else:
                                hours = str(difference.hours) + " hours "
                            if difference.minutes == 1:
                                minutes = "and 1 minute."
                            else:
                                minutes = "and " + str(difference.minutes) + " minutes."
                        else:
                            days = "0 days "
                            hours = "0 hours "
                            minutes = "and 0 minutes."
                        ransomware_results.append("Ransomware activity was prevented by security software within " + days + hours + minutes)
                    else:
                        ransomware_results.append("Ransomware activity was not prevented by security software.")

                if "detected by security and/or IT personnel" in ele['description']:
                    if ele['trigger'] == "Y":
                        if ele['time_start'] and ele['time_end']:
                            start = datetime.datetime.fromisoformat(ele['time_start'][:-1]).astimezone(timezone.utc)
                            end = datetime.datetime.fromisoformat(ele['time_end'][:-1]).astimezone(timezone.utc)
                            difference = relativedelta(end, start)

                            if difference.days == 1:
                                days = "1 day "
                            else:
                                days = str(difference.days) + " days "
                            if difference.hours == 1:
                                hours = "1 hour "
                            else:
                                hours = str(difference.hours) + " hours "
                            if difference.minutes == 1:
                                minutes = "and 1 minute."
                            else:
                                minutes = "and " + str(difference.minutes) + " minutes."
                        else:
                            days = "0 days "
                            hours = "0 hours "
                            minutes = "and 0 minutes."
                        ransomware_results.append("Ransomware activity was detected by security and/or IT personnel within " + days + hours + minutes)
                    else:
                        ransomware_results.append("Ransomware activity was not detected by security and/or IT personnel.")

                if "reported by end users" in ele['description']:
                    if ele['trigger'] == "Y":
                        if ele['time_start'] and ele['time_end']:
                            start = datetime.datetime.fromisoformat(ele['time_start'][:-1]).astimezone(timezone.utc)
                            end = datetime.datetime.fromisoformat(ele['time_end'][:-1]).astimezone(timezone.utc)
                            difference = relativedelta(end, start)

                            if difference.days == 1:
                                days = "1 day "
                            else:
                                days = str(difference.days) + " days "
                            if difference.hours == 1:
                                hours = "1 hour "
                            else:
                                hours = str(difference.hours) + " hours "
                            if difference.minutes == 1:
                                minutes = "and 1 minute."
                            else:
                                minutes = "and " + str(difference.minutes) + " minutes."
                        else:
                            days = "0 days "
                            hours = "0 hours "
                            minutes = "and 0 minutes."
                        ransomware_results.append("Ransomware activity was reported by end users within " + days + hours + minutes)
                    else:
                        ransomware_results.append("Ransomware activity was not reported by end users.")

        for line in ransomware_results:
            paragraph = ransomware_placeholder.add_paragraph()
            paragraph.level = 1
            run = paragraph.add_run()
            run.text = line
            paragraph.font.color.rgb = gray

    else:
        print("No ransomware data.")

    # insert data exfiltration slide
    if af.get_db_info(rva_info, 'dataexfil', 'NA') is not None:
        data_exfil_results = []
        
        for de in af.model_gen(rva_info, 'ptportal.dataexfil'):
            ele = de['fields']
            data = {"type": ele['datatype'], "protocol": ele['protocol'], "detection": ele['detection'], "prevention": ele['prevention']}
            data_exfil_results.append(data)

        data_exfil_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(data_exfil_layout)
        shapes = slide.shapes
        shapes.title.text = "DATA EXFILTRATION"

        # data exfil table
        rows = len(data_exfil_results) + 1
        cols = 4
        left = Inches(0.65)
        top = Inches(1.5)
        width = Inches(8.71)
        height = Inches(rows * 0.4)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        tbl = table._graphic_frame._element.graphic.graphicData.tbl
        tbl[0][-1].text = MediumStyle4Accent4

        # set column widths
        table.columns[0].width = Inches(3.2)
        table.columns[1].width = Inches(1.8)
        table.columns[2].width = Inches(1.85)
        table.columns[3].width = Inches(1.85)

        cell_text(table, 0, 0, "Data Type", "c", white)
        cell_text(table, 0, 1, "Protocol", "c", white)
        cell_text(table, 0, 2, "Detection", "c", white)
        cell_text(table, 0, 3, "Prevention", "c", white)

        for i in range(0, 4):
            table.cell(0, i).fill.solid()
            table.cell(0, i).fill.fore_color.rgb = blue
            table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(0, i).vertical_anchor = MSO_ANCHOR.MIDDLE

        for cnt, item in enumerate(data_exfil_results):
            cell_text(table, cnt + 1, 0, item['type'], color=gray)
            cell_text(table, cnt + 1, 1, item['protocol'], color=gray)

            det_cell = table.cell(cnt + 1, 2).text_frame.paragraphs[0]
            det_color = det_cell.add_run()
            det_color.text = ("• ")

            if item['detection'] == "D":
                det_color.font.color.rgb = low
                detection = det_cell.add_run()
                detection.text = "Detected"
                detection.font.color.rgb = gray
            else:
                det_color.font.color.rgb = crit
                detection = det_cell.add_run()
                detection.text = "Not Detected"
                detection.font.color.rgb = gray

            pre_cell = table.cell(cnt + 1, 3).text_frame.paragraphs[0]
            pre_color = pre_cell.add_run()
            pre_color.text = ("• ")

            if item['prevention'] == "B":
                pre_color.font.color.rgb = low
                prevention = pre_cell.add_run()
                prevention.text = "Blocked"
                prevention.font.color.rgb = gray
            else:
                pre_color.font.color.rgb = crit
                prevention = pre_cell.add_run()
                prevention.text = "Not Blocked"
                prevention.font.color.rgb = gray

            for i in range(0, 4):
                if i == 2 or i == 3:
                    table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[0].font.bold = True
                    table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[1].font.bold = True
                    table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                    table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[1].font.size = Pt(11)
                else:
                    table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[0].font.size = Pt(11)

                table.cell(cnt + 1, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                table.cell(cnt + 1, i).vertical_anchor = MSO_ANCHOR.MIDDLE
                table.cell(cnt + 1, i).fill.solid()
                table.cell(cnt + 1, i).fill.fore_color.rgb = white

    else:
        print("No data exfiltration data.")

    # insert payload testing slide
    if af.get_db_info(rva_info, 'payload', 'NA') is not None:
        payload_testing_results = []
        
        for pt in af.model_gen(rva_info, 'ptportal.payload'):
            ele = pt['fields']
            payload = {"description": ele['payload_description'], "protocol": ele['c2_protocol'], "host": ele['host_protection'], "border": ele['border_protection']}
            payload_testing_results.append(payload)

        chunks = [payload_testing_results[x:x+10] for x in range(0, len(payload_testing_results), 10)]

        for cnt, chunk in enumerate(chunks):
            payload_testing_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(payload_testing_layout)
            shapes = slide.shapes
            shapes.title.text = "PAYLOAD TESTING"

            # payload table
            rows = len(chunk) + 1
            cols = 4
            left = Inches(0.65)
            top = Inches(1.5)
            width = Inches(8.71)
            height = Inches(rows * 0.4)
            table = shapes.add_table(rows, cols, left, top, width, height).table
            tbl = table._graphic_frame._element.graphic.graphicData.tbl
            tbl[0][-1].text = MediumStyle4Accent4

            # set column widths
            table.columns[0].width = Inches(4)
            table.columns[1].width = Inches(1.3)
            table.columns[2].width = Inches(1.7)
            table.columns[3].width = Inches(1.7)

            cell_text(table, 0, 0, "Payload Description", "c", white)
            cell_text(table, 0, 1, "C2 Protocol", "c", white)
            cell_text(table, 0, 2, "Host Protection", "c", white)
            cell_text(table, 0, 3, "Border Protection", "c", white)

            for i in range(0, 4):
                table.cell(0, i).fill.solid()
                table.cell(0, i).fill.fore_color.rgb = blue
                table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(11)
                table.cell(0, i).vertical_anchor = MSO_ANCHOR.MIDDLE

            for cnt, item in enumerate(chunk):
                cell_text(table, cnt + 1, 0, item['description'], color=gray)
                cell_text(table, cnt + 1, 1, item['protocol'], color=gray)

                host_cell = table.cell(cnt + 1, 2).text_frame.paragraphs[0]
                host_color = host_cell.add_run()
                host_color.text = ("• ")

                if item['host'] == "B":
                    host_color.font.color.rgb = low
                    host_protection = host_cell.add_run()
                    host_protection.text = "Blocked"
                    host_protection.font.color.rgb = gray
                else:
                    host_color.font.color.rgb = crit
                    host_protection = host_cell.add_run()
                    host_protection.text = "Not Blocked"
                    host_protection.font.color.rgb = gray

                bord_cell = table.cell(cnt + 1, 3).text_frame.paragraphs[0]
                bord_color = bord_cell.add_run()
                bord_color.text = ("• ")

                if item['border'] == "B":
                    bord_color.font.color.rgb = low
                    border_protection = bord_cell.add_run()
                    border_protection.text = "Blocked"
                    border_protection.font.color.rgb = gray
                else:
                    bord_color.font.color.rgb = crit
                    border_protection = bord_cell.add_run()
                    border_protection.text = "Not Blocked"
                    border_protection.font.color.rgb = gray

                for i in range(0, 4):
                    if i == 2 or i == 3:
                        table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[0].font.bold = True
                        table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[1].font.bold = True
                        table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                        table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[1].font.size = Pt(11)
                    else:
                        table.cell(cnt + 1, i).text_frame.paragraphs[0].runs[0].font.size = Pt(11)

                    table.cell(cnt + 1, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    table.cell(cnt + 1, i).vertical_anchor = MSO_ANCHOR.MIDDLE
                    table.cell(cnt + 1, i).fill.solid()
                    table.cell(cnt + 1, i).fill.fore_color.rgb = white
    else:
        print("No payload testing results.")

    # insert phishing campaign slide
    if af.get_db_info(rva_info, 'campaign', 'NA') is not None:
        phishing_campaign_results = []

        for pc in af.model_gen(rva_info, 'ptportal.campaign'):
            ele = pc['fields']
            campaign = {"sent": ele['emails_sent'], "delivered": ele['emails_delivered'], "rate": ele['click_rate'], "total": ele['total_clicks'], "unique": ele['unique_clicks'], "first": ele['time_to_first_click'], "harvest": ele['creds_harvested'], "exploit": ele['number_exploited'], "length": ele['length_of_campaign']}
            phishing_campaign_results.append(campaign)

        for cnt, c in enumerate(phishing_campaign_results):
            phishing_campaign_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(phishing_campaign_layout)
            shapes = slide.shapes
            shapes.title.text = "PHISHING CAMPAIGN"

            # phishing campaign table
            rows = 10
            cols = 2
            left = Inches(0.65)
            top = Inches(1.5)
            width = Inches(8.71)
            height = Inches(rows * 0.4)
            table = shapes.add_table(rows, cols, left, top, width, height).table
            tbl = table._graphic_frame._element.graphic.graphicData.tbl
            tbl[0][-1].text = MediumStyle4Accent4

            # set column widths
            table.columns[0].width = Inches(5.3)
            table.columns[1].width = Inches(3.4)

            cell_text(table, 0, 0, "Campaign #" + str(cnt + 1), "c", color=white)
            table.cell(0, 0).merge(table.cell(0, 1))
            table.cell(0, 0).fill.solid()
            table.cell(0, 0).fill.fore_color.rgb = blue
            table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(0, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

            click_rate = round(float(c['rate']) * 100, 2)
            click_time = str(c['first']).split(" ")

            if len(click_time) > 1:
                if int(click_time[0]) == 1:
                    time_to_first_click = click_time[0] + " Day, " + click_time[1]
                else:
                    time_to_first_click = click_time[0] + " Days, " + click_time[1]
            else:
                time_to_first_click = str(c['first'])

            if int(c['length']) == 1:
                length = str(c['length']) + " Day"
            else:
                length = str(c['length']) + " Days"

            cell_text(table, 1, 0, "Emails Sent", color=gray)
            cell_text(table, 2, 0, "Emails Successfully Delivered", color=gray)
            cell_text(table, 3, 0, "Click Rate", color=gray)
            cell_text(table, 4, 0, "Total Clicks", color=gray)
            cell_text(table, 5, 0, "Unique Clicks", color=gray)
            cell_text(table, 6, 0, "Time to First Click (HH:MM:SS)", color=gray)
            cell_text(table, 7, 0, "Credentials Harvested", color=gray)
            cell_text(table, 8, 0, "Users Exploited", color=gray)
            cell_text(table, 9, 0, "Length of Campaign", color=gray)

            cell_text(table, 1, 1, str(c['sent']), color=gray)
            cell_text(table, 2, 1, str(c['delivered']), color=gray)
            cell_text(table, 3, 1, str(click_rate) + "%", color=gray)
            cell_text(table, 4, 1, str(c['total']), color=gray)
            cell_text(table, 5, 1, str(c['unique']), color=gray)
            cell_text(table, 6, 1, time_to_first_click, color=gray)
            cell_text(table, 7, 1, str(c['harvest']), color=gray)
            cell_text(table, 8, 1, str(c['exploit']), color=gray)
            cell_text(table, 9, 1, length, color=gray)

            for i in range(1, 10):
                for j in range(0, 2):
                    table.cell(i, j).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                    table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    table.cell(i, j).vertical_anchor = MSO_ANCHOR.MIDDLE
                    table.cell(i, j).fill.solid()
                    table.cell(i, j).fill.fore_color.rgb = white

    else:
        print("No phishing campaign data.")


def insert_attack_paths(prs, rva_info, ss_info, media_path):
    # ---- add Attack Path Section Slide
    title_only = prs.slide_layouts[11]
    slide = prs.slides.add_slide(title_only)
    slide.placeholders[13].text = "ATTACK PATHS"


    for cnt, path in enumerate(af.model_gen(rva_info, "ptportal.narrative")):
        ele = path['fields']
        name = ele['name'] + " " + str(ele['order'])

        attack_path_layout = prs.slide_layouts[15]
        slide = prs.slides.add_slide(attack_path_layout)
        shapes = slide.shapes
        shapes.title.text = name.upper()

        if ele['file']:
            dfile = media_path + ele['file']
            (x, y, w, h) = iu.get_screenshot_dimensions(dfile, "path")

            shapes.add_picture(dfile, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            slide.placeholders[12].text = ele['caption']


def insert_conclusion_slides(prs, rva_info):
    # ---- add Conclusion Section Slide
    title_only = prs.slide_layouts[11]
    slide = prs.slides.add_slide(title_only)
    slide.placeholders[13].text = "CONCLUSION"

    # insert Observations slide
    observation_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(observation_layout)
    title = slide.shapes.title
    title.text = "OBSERVATIONS"
    obs_placeholder = slide.placeholders[1].text_frame

    observations = af.get_db_info(rva_info, 'report.fields.observed_strengths', "").replace("\r", "\n")

    for cnt, bullet in enumerate(observations.split("\n")[::-1]):
        if cnt == 0:
            paragraph = obs_placeholder.paragraphs[0]
        else:
            paragraph = obs_placeholder.add_paragraph()
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = bullet
        paragraph.font.color.rgb = gray

    # insert Next Steps slide
    next_steps_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(next_steps_layout)
    title = slide.shapes.title
    title.text = "NEXT STEPS"
    ns_placeholder = slide.placeholders[1].text_frame

    poc_name = af.get_db_info(
        rva_info, "engagementmeta.fields.customer_POC_name", "Customer Name"
    )
    fed_name = af.get_db_info(
        rva_info, "engagementmeta.fields.team_lead_name", "Fed Lead Name"
    )
    customer_initials = af.get_db_info(
        rva_info, "engagementmeta.fields.customer_initials", "Stakeholder Initials"
    )

    paragraph = ns_placeholder.paragraphs[0]
    paragraph.level = 0
    run = paragraph.add_run()
    run.text = fed_name
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Additional analysis of assessment data"
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Send report draft to " + customer_initials
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Finalize report after review"
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 0
    run = paragraph.add_run()
    run.text = poc_name
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Review and validate findings"
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Provide high level report draft feedback"
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Create mitigation plans, as appropriate"
    paragraph.font.color.rgb = gray

    paragraph = ns_placeholder.add_paragraph()
    paragraph.level = 1
    run = paragraph.add_run()
    run.text = "Consider future work with CISA"
    paragraph.font.color.rgb = gray

    # insert Questions slide
    questions_layout = prs.slide_layouts[12]
    slide = prs.slides.add_slide(questions_layout)

    # insert final slide
    final_layout = prs.slide_layouts[17]
    slide = prs.slides.add_slide(final_layout)


def generate_ptp_slides(template, output, draft, json, media):
    """Generates an outbrief for the current assessment.

    Args:
        template (string): Path to the docx template that will be used to generate the report.
        output (string): Name of the file that will be saved and returned.
        draft (boolean): Marks the docx document with a draft watermark.
        json (string): Path to the json file with the assessment data.
        media (string): Path to the media folder that contains the assessment screenshots.
    """

    # ---- find all screenshot information
    rva_info = af.load_rva_info(json)
    ss_info = af.build_screenshot_info(rva_info)

    rep_fields = af.get_db_info(rva_info, "report.fields", "keyNA")
    report_type = rep_fields["report_type"]
    ip_ext_scan = rep_fields["external_scanned"]
    ip_ext_disc = rep_fields["external_discovered"]
    ip_int_scan = rep_fields["internal_scanned"]
    ip_int_disc = rep_fields["internal_discovered"]

    # ---- open the powerpoint template
    prs = pptx.Presentation(template)

    insert_title_slide(prs, report_type, rva_info, draft)
    insert_notice_slide(prs, report_type)
    insert_agenda_slide(prs, report_type)
    insert_logistics(prs, report_type, rva_info)
    insert_scope_slide(prs, report_type, rva_info, ip_ext_scan, ip_ext_disc, ip_int_scan, ip_int_disc)
    insert_goals_slide(prs, report_type)

    insert_findings_slides(prs, report_type, rva_info, ss_info, media)
    insert_services_slides(prs, rva_info)
    insert_attack_paths(prs, rva_info, ss_info, media)

    insert_conclusion_slides(prs, rva_info)

    #if report_type == "RPT":
    #    insert_OSINT_slide(prs, report_type, rva_info)

    prs.save(output)


def main():
    description = "Generate RVA briefing slide deck"
    parser = argparse.ArgumentParser(description=description)

    parser.add_argument("TEMPLATE", help="Briefing slide deck template.")
    parser.add_argument(
        "-o",
        "--output_file",
        action="store",
        default="RVA_OutBrief_" + tstamp + ".pptx",
        help="Out briefing file name",
    )
    parser.add_argument(
        "-d", "--draft", action="store_true", help="Label report as a draft"
    )
    parser.add_argument("-j", "--json_file", action="store", required=True)
    parser.add_argument(
        "-m",
        "--media_path",
        action="store",
        default="./",
        help="Location of screenshots, etc.",
    )
    args = parser.parse_args()

    generate_ptp_slides(
        args.TEMPLATE, args.output_file, args.draft, args.json_file, args.media_path
    )


if __name__ == "__main__":
    main()
